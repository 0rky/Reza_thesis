from __future__ import annotations

import argparse
import csv
from pathlib import Path
from typing import Optional

from pptx import Presentation
from pptx.enum.shapes import MSO_SHAPE_TYPE
from pptx.util import Pt


# -----------------------------
# Helpers
# -----------------------------

def read_csv_rows(path: Path):
    with path.open(newline="") as f:
        return list(csv.DictReader(f))


def to_float(v: Optional[str]) -> Optional[float]:
    if v is None:
        return None
    s = str(v).strip()
    if not s:
        return None
    try:
        return float(s)
    except Exception:
        return None


def find_slide_by_title(prs: Presentation, title_text: str):
    target = title_text.strip().lower()
    for slide in prs.slides:
        for shape in slide.shapes:
            if getattr(shape, "has_text_frame", False):
                text = (shape.text or "").strip().lower()
                if text == target:
                    return slide
    return None


def find_text_shapes(slide):
    return [s for s in slide.shapes if getattr(s, "has_text_frame", False)]


def pick_largest_text_shape_excluding_title(slide, title_contains: str):
    title_contains = title_contains.lower()
    candidates = []
    for s in find_text_shapes(slide):
        txt = (s.text or "").strip().lower()
        if not txt:
            continue
        if title_contains in txt:
            continue
        area = int(s.width) * int(s.height)
        candidates.append((area, s))
    candidates.sort(key=lambda x: x[0], reverse=True)
    return candidates[0][1] if candidates else None


def set_shape_text(shape, text: str, font_size: Optional[int] = None):
    tf = shape.text_frame
    tf.clear()
    paragraphs = text.split("\n")
    first = True
    for line in paragraphs:
        if first:
            p = tf.paragraphs[0]
            first = False
        else:
            p = tf.add_paragraph()
        p.text = line
        if font_size is not None:
            for run in p.runs:
                run.font.size = Pt(font_size)


def replace_picture_keep_frame(slide, old_pic_shape, image_path: Path):
    if not image_path.exists():
        print(f"[WARN] Missing image, skipping: {image_path}")
        return None
    left, top, width, height = old_pic_shape.left, old_pic_shape.top, old_pic_shape.width, old_pic_shape.height
    el = old_pic_shape._element
    spTree = el.getparent()
    spTree.remove(el)
    return slide.shapes.add_picture(str(image_path), left, top, width=width, height=height)


def list_pictures(slide):
    return [s for s in slide.shapes if s.shape_type == MSO_SHAPE_TYPE.PICTURE]


def update_table(slide, rows):
    table_shape = None
    for s in slide.shapes:
        if s.shape_type == MSO_SHAPE_TYPE.TABLE:
            table_shape = s
            break
    if table_shape is None:
        print(f"[WARN] No table found on slide: {slide}")
        return

    table = table_shape.table
    n_rows = len(table.rows)
    n_cols = len(table.columns)

    if n_cols != 5:
        print(f"[WARN] Final comparison table has {n_cols} columns, expected 5. Skipping table update.")
        return

    header = [
        "Method",
        "Accuracy",
        "Loss",
        "Compression / Sparsity",
        "Main takeaway",
    ]
    for c, val in enumerate(header):
        table.cell(0, c).text = val

    method_name_map = {
        "baseline": "Baseline",
        "mixed_precision": "Mixed Precision",
        "quantization_dynamic": "Dynamic Quantization",
        "qat_converted": "QAT Converted",
        "pruning_pruned": "Weight Pruning",
        "structured_pruned": "Structured Pruning",
    }

    takeaway_map = {
        "baseline": "Strong reference after full 10-epoch run",
        "mixed_precision": "Lower GPU memory, but not best overall",
        "quantization_dynamic": "Accuracy preserved, compression still weak",
        "qat_converted": "Best overall: strongest accuracy/compression trade-off",
        "pruning_pruned": "Accuracy preserved, but dense storage prevents real file compression",
        "structured_pruned": "Recovery fine-tuning fixed the earlier collapse",
    }

    def compression_text(row):
        case = row["case"]
        reduction = to_float(row.get("size_reduction_percent"))
        sparsity = to_float(row.get("sparsity_percent"))
        structured_sparsity = to_float(row.get("structured_sparsity_percent"))

        if case == "baseline":
            return "Reference"
        if case == "mixed_precision":
            return "Training-time memory reduction"
        if case == "quantization_dynamic":
            return f"{reduction:.2f}% size reduction" if reduction is not None else "Quantized inference"
        if case == "qat_converted":
            return f"{reduction:.2f}% size reduction" if reduction is not None else "Quantization-aware training"
        if case == "pruning_pruned":
            return f"{sparsity:.2f}% sparsity" if sparsity is not None else "Pruned"
        if case == "structured_pruned":
            return f"{structured_sparsity:.2f}% structured sparsity" if structured_sparsity is not None else "Structured pruning"
        return ""

    max_data_rows = n_rows - 1
    usable_rows = rows[:max_data_rows]

    for r, row in enumerate(usable_rows, start=1):
        acc = to_float(row.get("accuracy_percent"))
        loss = to_float(row.get("cross_entropy_loss"))

        vals = [
            method_name_map.get(row["case"], row["case"]),
            f"{acc:.2f}%" if acc is not None else "",
            f"{loss:.6f}" if loss is not None else "",
            compression_text(row),
            takeaway_map.get(row["case"], ""),
        ]

        for c, val in enumerate(vals):
            table.cell(r, c).text = val

    print("[OK] Final comparison table updated with 5-column layout")


# -----------------------------
# Content builders
# -----------------------------

def get_case(rows, name: str):
    for row in rows:
        if row["case"] == name:
            return row
    raise KeyError(name)


def build_experiment_texts(rows, enriched_rows):
    baseline = get_case(rows, "baseline")
    mixed = get_case(rows, "mixed_precision")
    qfloat = get_case(rows, "quantization_float")
    qdyn = get_case(rows, "quantization_dynamic")
    qat_pre = get_case(rows, "qat_preconvert")
    qat_conv = get_case(rows, "qat_converted")
    pdense = get_case(rows, "pruning_dense")
    ppruned = get_case(rows, "pruning_pruned")
    sdense = get_case(rows, "structured_dense")
    spruned = get_case(rows, "structured_pruned")

    mixed_e = get_case(enriched_rows, "mixed_precision")
    qfloat_e = get_case(enriched_rows, "quantization_float")
    qat_pre_e = get_case(enriched_rows, "qat_preconvert")
    spruned_e = get_case(enriched_rows, "structured_pruned")

    texts = {}

    texts["Baseline: setup and implementation"] = (
        "Dataset: FashionMNIST\n"
        "Input resized to 224 x 224\n"
        "Batch size: 64\n"
        "Epochs: 10\n"
        "Optimizer: Adam\n"
        "Loss: CrossEntropyLoss\n"
        "Architecture: multiple Conv2d layers, pooling, dropout, flatten, and fully connected classification head\n"
        "Evaluation and profiling added: confusion matrix, classification report, metrics export, accuracy/loss curves, runtime summary, epoch profiling, and layer-wise profiling"
    )

    texts["Baseline: experiment result"] = (
        f"Accuracy: {to_float(baseline['accuracy_percent']):.2f}%\n"
        f"Loss: {to_float(baseline['cross_entropy_loss']):.6f}\n"
        "Performance improved after the full 10-epoch run and now provides a stronger reference point for all comparisons\n"
        "Good performance remains visible on distinctive classes such as Trouser, Bag, Sandal, and Ankle boot\n"
        "Main confusion is still concentrated among similar upper-body classes such as T-shirt/top, Pullover, Coat, and Shirt"
    )

    texts["Mixed Precision: concept, purpose, strengths, limitations"] = (
        "Mixed precision training uses more than one numeric precision during training, typically FP16 and FP32\n\n"
        "Purpose: reduce GPU memory pressure and potentially improve training throughput\n"
        "Advantages: lower GPU memory usage, possible speed improvement, relatively clean integration into the training loop\n"
        "Disadvantages: hardware-dependent benefit, possible numerical instability, and no guarantee of better final accuracy\n\n"
        "Important clarification: mixed precision is mainly a training-time technique\n"
        "It does not automatically produce a smaller saved model file\n"
        "A key follow-up question is whether precision selection can be controlled manually rather than relying only on automatic autocast behavior"
    )

    texts["Mixed Precision: implementation in this code"] = (
        "The model architecture was not redesigned\n"
        "The main computation path was modified using autocast and GradScaler\n"
        "The training loop changed around forward pass, loss computation, and backward scaling\n"
        "Data pipeline and optimizer remained the same\n\n"
        "Code-level summary\n"
        "- training forward pass moved into mixed precision context\n"
        "- loss scaling added for stable optimization\n"
        "- optimizer and data pipeline remained unchanged\n"
        "- profiling was added to measure epoch time, GPU memory, RSS RAM, and layer-wise behavior\n\n"
        "Professor follow-up to address in discussion:\n"
        "Can precision handling be controlled manually instead of leaving all casting decisions to automatic selection?"
    )

    texts["Mixed Precision: experiment result"] = (
        f"Accuracy: {to_float(mixed['accuracy_percent']):.2f}%\n"
        f"Loss: {to_float(mixed['cross_entropy_loss']):.6f}\n"
        f"Average train epoch time: {to_float(mixed_e['avg_train_epoch_seconds']):.2f} seconds\n"
        f"Max GPU allocated: {to_float(mixed_e['max_gpu_allocated_mb']):.2f} MB\n\n"
        "Compared with baseline, mixed precision reduced peak GPU allocation substantially, but final accuracy was slightly lower\n"
        "In this Tesla P100 setup, the practical advantage came mostly from memory reduction rather than a dramatic accuracy or runtime gain\n"
        "Conclusion in this setup: feasible and useful for memory pressure, but not the strongest overall method"
    )

    texts["Dynamic Quantization: Experiment result"] = (
        f"Float model accuracy: {to_float(qfloat['accuracy_percent']):.2f}%\n"
        f"Float model loss: {to_float(qfloat['cross_entropy_loss']):.6f}\n"
        f"Dynamic quantized model accuracy: {to_float(qdyn['accuracy_percent']):.2f}%\n"
        f"Dynamic quantized model loss: {to_float(qdyn['cross_entropy_loss']):.6f}\n"
        f"Size reduction: {to_float(qdyn['size_reduction_percent']):.2f}%\n\n"
        "Accuracy was essentially preserved, but compression remained negligible\n"
        "This result supports the point that the chosen dynamic quantization path here mainly affected deployability slightly, not training speed\n"
        "Important conceptual point: this quantization path is primarily an inference/deployment story, not a training-time acceleration story"
    )

    texts["QAT: concept, purpose, strengths, limitations"] = (
        "QAT stands for Quantization-Aware Training\n"
        "Instead of quantizing only after training, the model is trained while simulating quantization effects\n\n"
        "Purpose: preserve or improve final quantized accuracy while still achieving real compression\n"
        "Advantages: often preserves quantized performance better than naive post-training quantization and can yield real deployable compression\n"
        "Disadvantages: more code changes, more setup complexity, and explicit quantization preparation logic is required\n\n"
        "Important clarification: QAT is training-time preparation for a better final quantized inference model\n"
        "This is why it bridges the training/inference gap more directly than simple post-training dynamic quantization"
    )

    texts["QAT: experiment result"] = (
        f"QAT pre-convert accuracy: {to_float(qat_pre['accuracy_percent']):.2f}%\n"
        f"QAT pre-convert loss: {to_float(qat_pre['cross_entropy_loss']):.6f}\n"
        f"QAT converted accuracy: {to_float(qat_conv['accuracy_percent']):.2f}%\n"
        f"QAT converted loss: {to_float(qat_conv['cross_entropy_loss']):.6f}\n"
        f"Size reduction: {to_float(qat_conv['size_reduction_percent']):.2f}%\n"
        f"Model size dropped from {to_float(qat_pre['state_dict_size_mb']):.4f} MB to {to_float(qat_conv['state_dict_size_mb']):.4f} MB\n\n"
        "This is now the strongest result in the study\n"
        "QAT achieved the best accuracy while also producing a substantially smaller deployable model\n"
        "This makes QAT the best overall trade-off between performance and compression in the current experiments"
    )

    texts["Weight Pruning: experiment result"] = (
        f"Dense model accuracy: {to_float(pdense['accuracy_percent']):.2f}%\n"
        f"Dense model loss: {to_float(pdense['cross_entropy_loss']):.6f}\n\n"
        f"Pruned model accuracy: {to_float(ppruned['accuracy_percent']):.2f}%\n"
        f"Pruned model loss: {to_float(ppruned['cross_entropy_loss']):.6f}\n\n"
        f"Global sparsity: {to_float(ppruned['sparsity_percent']):.2f}%\n"
        f"File size reduction: {to_float(ppruned['size_reduction_percent']):.2f}%\n\n"
        "Accuracy was preserved and even improved slightly after pruning\n"
        "However, saved model size did not decrease in a meaningful way because the model was still stored in dense format\n"
        "This directly supports the professor's point that sparsity alone does not guarantee storage or hardware benefit"
    )

    texts["Structured Pruning: experiment result"] = (
        f"Dense model accuracy: {to_float(sdense['accuracy_percent']):.2f}%\n"
        f"Dense model loss: {to_float(sdense['cross_entropy_loss']):.6f}\n"
        f"Structured-pruned model accuracy: {to_float(spruned['accuracy_percent']):.2f}%\n"
        f"Structured-pruned model loss: {to_float(spruned['cross_entropy_loss']):.6f}\n"
        f"Zero-weight sparsity: {to_float(spruned['zero_weight_sparsity_percent']):.2f}%\n"
        f"Structured sparsity: {to_float(spruned['structured_sparsity_percent']):.2f}%\n"
        f"File size reduction: {to_float(spruned['size_reduction_percent']):.2f}%\n\n"
        "This result changed substantially after fixing the procedure and adding short fine-tuning after structured pruning\n"
        "Performance no longer collapses\n"
        "Even with only about 9.89% structured sparsity, the model recovered and outperformed the structured-dense case\n"
        "Takeaway: naive structured pruning was too destructive before, but structured pruning plus recovery/fine-tuning is much more reasonable"
    )

    texts["Final conclusion"] = (
        "Main findings\n"
        f"- Baseline after the full run reached {to_float(baseline['accuracy_percent']):.2f}% accuracy\n"
        f"- Mixed precision reduced GPU memory strongly but did not beat the best final accuracy\n"
        f"- Dynamic quantization preserved accuracy but reduced size by only {to_float(qdyn['size_reduction_percent']):.2f}%\n"
        f"- Weight pruning preserved accuracy at {to_float(ppruned['sparsity_percent']):.2f}% sparsity, but dense storage prevented meaningful file compression\n"
        f"- Structured pruning no longer collapsed after adding recovery fine-tuning, and reached {to_float(spruned['accuracy_percent']):.2f}%\n"
        f"- QAT gave the strongest overall result: {to_float(qat_pre['accuracy_percent']):.2f}% pre-convert accuracy and {to_float(qat_conv['size_reduction_percent']):.2f}% size reduction after conversion\n\n"
        "Final conclusion\n"
        "Under this baseline model, dataset, and Tesla P100 hardware, QAT was the most effective method overall\n"
        "The next step is not more polishing on this toy model; it is to carry the lessons forward into larger models and the real thesis dataset"
    )

    texts["Professor Questions"] = (
        "Professor-driven clarifications to address explicitly\n\n"
        "1. Mixed precision vs quantization\n"
        "- mixed precision is mainly a training-time compute/memory technique\n"
        "- dynamic quantization is mainly an inference/deployment technique\n"
        "- QAT links training to a final quantized inference model\n\n"
        "2. Manual control question\n"
        "- current code uses automatic autocast behavior\n"
        "- next improvement is to test explicit/manual precision control where appropriate instead of relying only on automatic casting\n\n"
        "3. Memory interpretation question\n"
        "- GPU memory is the main memory of interest here, not generic RAM language\n"
        "- throughput should be described concretely as work processed per unit time, not vaguely as faster pipeline behavior\n\n"
        "4. Sparsity interpretation question\n"
        "- zero weights do not guarantee smaller files or faster execution if storage and kernels remain dense\n\n"
        "5. Layer-wise profiling question\n"
        "- layer-wise profiling has now been added so later discussion can move from vague claims to measured hotspots"
    )

    texts["Layer-wise profiling"] = (
        f"Layer-wise profiling was added for baseline and mixed precision\n"
        f"Baseline peak GPU allocation: 438.56 MB\n"
        f"Mixed precision peak GPU allocation: {to_float(mixed_e['max_gpu_allocated_mb']):.2f} MB\n"
        f"Reduction in peak GPU allocation: {438.56 - to_float(mixed_e['max_gpu_allocated_mb']):.2f} MB\n\n"
        "This profiling is intended to answer the professor's request to move beyond final accuracy only\n"
        "The new instrumentation captures per-epoch timing, GPU allocation, reserved memory, RSS RAM, and per-layer forward-pass measurements\n"
        "This makes it possible to discuss where memory expands/contracts and which layers dominate runtime or memory behavior"
    )

    return texts


# -----------------------------
# Main updater
# -----------------------------

def main():
    parser = argparse.ArgumentParser()
    parser.add_argument("--input", required=True)
    parser.add_argument("--output", required=True)
    parser.add_argument("--repo-root", required=True)
    args = parser.parse_args()

    input_pptx = Path(args.input)
    output_pptx = Path(args.output)
    repo = Path(args.repo_root)

    prs = Presentation(str(input_pptx))

    rows = read_csv_rows(repo / "reports" / "comparison_all_experiments.csv")
    enriched_rows = read_csv_rows(repo / "reports" / "comparison_all_experiments_enriched.csv")
    texts = build_experiment_texts(rows, enriched_rows)

    for title, new_text in texts.items():
        slide = find_slide_by_title(prs, title)
        if slide is None:
            print(f"[WARN] Slide not found by title: {title}")
            continue
        target = pick_largest_text_shape_excluding_title(slide, title)
        if target is None:
            print(f"[WARN] No body text box found for slide: {title}")
            continue
        set_shape_text(target, new_text, font_size=16)
        print(f"[OK] Updated text on slide: {title}")

    image_map = {
        "Baseline: experiment result": [
            repo / "baseline" / "outputs" / "baseline_confusion_matrix.png",
            repo / "baseline" / "outputs" / "baseline_loss_curve.png",
            repo / "baseline" / "outputs" / "baseline_accuracy_curve.png",
        ],
        "Mixed Precision: experiment result": [
            repo / "mixed_precision" / "outputs" / "mixed_precision_confusion_matrix.png",
            repo / "mixed_precision" / "outputs" / "mixed_precision_loss_curve.png",
            repo / "mixed_precision" / "outputs" / "mixed_precision_accuracy_curve.png",
            repo / "reports" / "plots_enriched" / "max_gpu_allocated_comparison.png",
        ],
        "Dynamic Quantization: Experiment result": [
            repo / "quantization" / "outputs" / "quantization_dynamic_confusion_matrix.png",
        ],
        "QAT: experiment result": [
            repo / "qat" / "outputs" / "qat_preconvert_confusion_matrix.png",
            repo / "qat" / "outputs" / "qat_converted_confusion_matrix.png",
        ],
        "Weight Pruning: experiment result": [
            repo / "pruning" / "outputs" / "pruning_dense_confusion_matrix.png",
            repo / "pruning" / "outputs" / "pruning_pruned_confusion_matrix.png",
        ],
        "Structured Pruning: experiment result": [
            repo / "structured_pruning" / "outputs" / "structured_dense_confusion_matrix.png",
            repo / "structured_pruning" / "outputs" / "structured_pruned_confusion_matrix.png",
        ],
    }

    for title, new_images in image_map.items():
        slide = find_slide_by_title(prs, title)
        if slide is None:
            continue
        pics = list_pictures(slide)
        for old_pic, new_img in zip(pics, new_images):
            replace_picture_keep_frame(slide, old_pic, new_img)
        print(f"[OK] Replaced images on slide: {title}")

    final_comp_slide = find_slide_by_title(prs, "Final comparison")
    if final_comp_slide is not None:
        selected_cases = [
            get_case(enriched_rows, "baseline"),
            get_case(enriched_rows, "mixed_precision"),
            get_case(enriched_rows, "quantization_dynamic"),
            get_case(enriched_rows, "qat_converted"),
            get_case(enriched_rows, "pruning_pruned"),
            get_case(enriched_rows, "structured_pruned"),
        ]
        update_table(final_comp_slide, selected_cases)
        print("[OK] Updated final comparison table")

    code_slide = find_slide_by_title(prs, "Code-level summary of modifications")
    if code_slide is not None:
        table_shape = None
        for s in code_slide.shapes:
            if s.shape_type == MSO_SHAPE_TYPE.TABLE:
                table_shape = s
                break
        if table_shape is not None:
            table = table_shape.table
            headers = ["Area", "What changed", "Why it matters"]
            for i, h in enumerate(headers[: len(table.columns)]):
                table.cell(0, i).text = h
            rows_content = [
                ("Training scripts", "All six cases now read EPOCHS from environment and support the automated 10-epoch run script", "Ensures fair and repeatable comparison"),
                ("Profiling", "Added runtime summary, epoch profiling, and layer-wise profiling", "Allows memory/time interpretation instead of accuracy-only reporting"),
                ("Reports", "Added aggregated CSV/Markdown reports and enriched comparison reports", "Centralizes results across methods"),
                ("Plots", "Added summary plots for accuracy, loss, size, reduction, timing, and GPU memory", "Makes cross-method comparison presentation-ready"),
                ("Structured pruning", "Added safer pruning logic plus short fine-tuning recovery", "Prevents the earlier collapse and yields a much stronger result"),
                ("Automation", "Added run_all_10ep_and_reports.sh", "Runs all cases sequentially and rebuilds reports automatically"),
            ]
            for r, rowvals in enumerate(rows_content, start=1):
                if r >= len(table.rows):
                    break
                for c, val in enumerate(rowvals[: len(table.columns)]):
                    table.cell(r, c).text = val
            print("[OK] Updated code summary table")

    best_slide = find_slide_by_title(prs, "Best and worst methods in this study")
    if best_slide is not None:
        body = None
        for s in best_slide.shapes:
            if getattr(s, "has_text_frame", False) and "Best and worst methods" not in (s.text or ""):
                body = s
                break
        if body is None:
            for s in best_slide.placeholders:
                if getattr(s, "has_text_frame", False) and s.placeholder_format.idx != 0:
                    body = s
                    break
        if body is not None:
            set_shape_text(
                body,
                "Best overall method\n"
                "- QAT was the strongest method overall\n"
                "- Best accuracy: 90.80% pre-convert, 90.79% converted\n"
                "- Best real compression: 74.75% size reduction after conversion\n\n"
                "Weakest result in the earlier version\n"
                "- structured pruning originally collapsed before recovery fine-tuning was added\n\n"
                "Updated interpretation\n"
                "- after the fix, structured pruning no longer belongs in the failure category\n"
                "- the real weak point now is that unstructured pruning and dynamic quantization gave little or no meaningful file-size reduction in this dense-storage setup",
                font_size=20,
            )
            print("[OK] Updated best/worst slide")

    blank_layout = prs.slide_layouts[6]
    for appendix_title, appendix_body, appendix_image in [
        (
            "Professor Questions",
            texts["Professor Questions"],
            repo / "reports" / "plots_enriched" / "accuracy_vs_gpu_memory.png",
        ),
        (
            "Layer-wise profiling",
            texts["Layer-wise profiling"],
            repo / "reports" / "plots_enriched" / "avg_train_epoch_time_comparison.png",
        ),
    ]:
        slide = prs.slides.add_slide(blank_layout)
        title_box = slide.shapes.add_textbox(Pt(30), Pt(20), Pt(600), Pt(40))
        set_shape_text(title_box, appendix_title, font_size=24)
        body_box = slide.shapes.add_textbox(Pt(30), Pt(70), Pt(420), Pt(420))
        set_shape_text(body_box, appendix_body, font_size=16)
        if appendix_image.exists():
            slide.shapes.add_picture(str(appendix_image), Pt(470), Pt(90), width=Pt(420), height=Pt(300))
        print(f"[OK] Added appendix slide: {appendix_title}")

    prs.save(str(output_pptx))
    print(f"[DONE] Saved updated presentation to: {output_pptx}")


if __name__ == "__main__":
    main()
